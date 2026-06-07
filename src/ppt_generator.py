import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE


GROUNDED_PATH = "outputs/grounded_slides.json"
OUTPUT_PPT = "outputs/paperbridge_grounded_slides.pptx"


# ---------- 基础样式 ----------

SLIDE_W = 13.333
SLIDE_H = 7.5

COLOR_BG = RGBColor(248, 250, 252)
COLOR_DARK = RGBColor(30, 41, 59)
COLOR_TEXT = RGBColor(51, 65, 85)
COLOR_MUTED = RGBColor(100, 116, 139)
COLOR_BLUE = RGBColor(37, 99, 235)
COLOR_LIGHT_BLUE = RGBColor(219, 234, 254)
COLOR_GREEN = RGBColor(22, 163, 74)
COLOR_LIGHT_GREEN = RGBColor(220, 252, 231)
COLOR_ORANGE = RGBColor(234, 88, 12)
COLOR_LIGHT_ORANGE = RGBColor(255, 237, 213)
COLOR_PURPLE = RGBColor(126, 34, 206)
COLOR_LIGHT_PURPLE = RGBColor(243, 232, 255)
COLOR_GRAY = RGBColor(226, 232, 240)
COLOR_WHITE = RGBColor(255, 255, 255)


def truncate(text, max_len):
    text = str(text).strip().replace("\n", " ")
    if len(text) <= max_len:
        return text
    return text[:max_len - 1] + "…"


def set_run_font(run, size=18, bold=False, color=COLOR_TEXT):
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def set_shape_text(shape, text, font_size=16, bold=False, color=COLOR_TEXT, align=None):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)

    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = str(text)
    set_run_font(run, size=font_size, bold=bold, color=color)


def add_bg(slide):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(SLIDE_W),
        Inches(SLIDE_H),
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG
    bg.line.fill.background()


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.25), Inches(12.2), Inches(0.65))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run_font(run, size=28, bold=True, color=COLOR_DARK)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(0.92), Inches(12.0), Inches(0.38))
        tf = sub.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        set_run_font(run, size=13, color=COLOR_MUTED)


def add_label(slide, text, x, y, w, h, color=COLOR_BLUE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    set_shape_text(shape, text, font_size=12, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    return shape


def add_card(slide, title, body, x, y, w, h, fill=COLOR_WHITE, border=COLOR_GRAY,
             title_color=COLOR_DARK, body_size=13):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.1)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(title)
    set_run_font(run, size=15, bold=True, color=title_color)

    if body:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        run2 = p2.add_run()
        run2.text = str(body)
        set_run_font(run2, size=body_size, color=COLOR_TEXT)

    return shape


def add_bullet_box(slide, title, items, x, y, w, h, max_item_len=105):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_WHITE
    shape.line.color.rgb = COLOR_GRAY
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.14)
    tf.margin_bottom = Inches(0.12)

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run_font(r, size=16, bold=True, color=COLOR_DARK)

    for item in items:
        p = tf.add_paragraph()
        p.level = 0
        p.space_before = Pt(5)
        r = p.add_run()
        r.text = "• " + truncate(item, max_item_len)
        set_run_font(r, size=14, color=COLOR_TEXT)

    return shape


def add_evidence_box(slide, evidence, x, y, w, h, max_items=4):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(241, 245, 249)
    shape.line.color.rgb = RGBColor(203, 213, 225)
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.10)
    tf.margin_bottom = Inches(0.08)

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Evidence from the paper"
    set_run_font(r, size=12.5, bold=True, color=COLOR_BLUE)

    if not evidence:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = "No evidence paragraphs selected."
        set_run_font(r, size=10.5, color=COLOR_MUTED)
        return shape

    # 小高度的 evidence box 用紧凑显示，避免思维导图/流程图底部文字溢出
    compact_mode = h < 1.3

    if compact_mode:
        show_count = min(max_items, 1, len(evidence))
    else:
        show_count = min(max_items, 3, len(evidence))

    for ev in evidence[:show_count]:
        pid = ev.get("paragraph_id", "?")
        section = ev.get("section_guess", "Other")
        page_start = ev.get("page_start", "?")
        page_end = ev.get("page_end", page_start)

        if page_start == page_end:
            loc = f"Page {page_start}"
        else:
            loc = f"Page {page_start}-{page_end}"

        excerpt = truncate(ev.get("source_excerpt", ""), 78 if not compact_mode else 95)
        summary = truncate(ev.get("summary_sentence", ""), 88 if not compact_mode else 70)

        p = tf.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run()

        if compact_mode:
            r.text = f"[P{pid} | {loc} | {section}] Original: {excerpt}"
            set_run_font(r, size=7.8, color=COLOR_TEXT)
        else:
            r.text = (
                f"[P{pid} | {loc} | {section}]\n"
                f"Original: {excerpt}\n"
                f"Summary: {summary}"
            )
            set_run_font(r, size=8.6, color=COLOR_TEXT)

    if len(evidence) > show_count and not compact_mode:
        p = tf.add_paragraph()
        p.space_before = Pt(2)
        r = p.add_run()
        r.text = f"... and {len(evidence) - show_count} more evidence blocks"
        set_run_font(r, size=8.2, color=COLOR_MUTED)

    return shape



def add_terms_box(slide, terms, x, y, w, h):
    if not terms:
        add_card(
            slide,
            "Key terms for beginners",
            "No key terms selected for this slide.",
            x, y, w, h,
            fill=COLOR_WHITE,
            border=COLOR_GRAY,
            title_color=COLOR_PURPLE,
            body_size=12,
        )
        return

    lines = []
    for item in terms[:4]:
        term = truncate(item.get("term", ""), 24)
        explanation = truncate(item.get("explanation", ""), 58)
        if term and explanation:
            lines.append(f"{term}: {explanation}")

    add_bullet_box(
        slide,
        "Key terms for beginners",
        lines if lines else ["No key terms selected for this slide."],
        x, y, w, h,
        max_item_len=95,
    )

def add_connector(slide, x1, y1, x2, y2, color=COLOR_MUTED):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(1.3)
    return line


def get_slide_by_no(slides, no):
    for s in slides:
        if int(s.get("slide_no", -1)) == no:
            return s
    return None


# ---------- 特殊页面：思维导图 ----------

def flatten_mind_children(node):
    return node.get("children", []) if isinstance(node, dict) else []


def collect_child_names(node, max_items=4):
    """
    把 mind_map 中某个一级节点下面的内容压成几条短 bullet。
    只取两层，避免 PPT 过密。
    """
    if not isinstance(node, dict):
        return []

    lines = []
    for child in node.get("children", []):
        name = child.get("name", "")
        if name:
            lines.append(name)

        # 如果还有二级子节点，挑最重要的一两个
        for sub in child.get("children", [])[:2]:
            sub_name = sub.get("name", "")
            if sub_name:
                lines.append(sub_name)

        if len(lines) >= max_items:
            break

    return lines[:max_items]


def add_map_column(slide, index, title, items, x, y, w, h, fill, border):
    """
    一列学习地图卡片：顶部编号 + 标题 + 3~4条短解释。
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(1.3)

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.10)
    tf.margin_bottom = Inches(0.08)

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"{index}. {title}"
    set_run_font(r, size=14.5, bold=True, color=border)

    for item in items[:4]:
        p = tf.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = "• " + truncate(item, 34)
        set_run_font(r, size=10.3, color=COLOR_TEXT)

    return shape


def draw_mind_map(slide, mind_map):
    """
    改进版：不再使用中心放射状连线。
    使用从左到右的 learning map：
    Prerequisite -> Problem -> Method -> Experiments -> Takeaways
    """
    children = flatten_mind_children(mind_map)

    # 按名称找对应节点，避免模型输出顺序变化
    def find_node(keyword):
        for child in children:
            name = child.get("name", "").lower()
            if keyword.lower() in name:
                return child
        return {}

    nodes = [
        ("Prerequisite", "Prerequisite Knowledge", find_node("Prerequisite"), COLOR_LIGHT_GREEN, COLOR_GREEN),
        ("Problem", "Research Problem", find_node("Problem"), COLOR_LIGHT_ORANGE, COLOR_ORANGE),
        ("Method", "Core Method", find_node("Method"), COLOR_LIGHT_PURPLE, COLOR_PURPLE),
        ("Experiments", "Experiments", find_node("Experiments"), COLOR_LIGHT_BLUE, COLOR_BLUE),
        ("Takeaways", "Takeaways", find_node("Takeaways"), RGBColor(254, 249, 195), RGBColor(202, 138, 4)),
    ]

    # 顶部核心主题
    root_name = truncate(mind_map.get("root", "Core Paper"), 95)
    add_card(
        slide,
        "Core idea of the paper",
        root_name,
        1.0,
        1.18,
        11.35,
        0.72,
        fill=COLOR_WHITE,
        border=COLOR_BLUE,
        title_color=COLOR_BLUE,
        body_size=12.5,
    )

    # 五列横向学习地图
    start_x = 0.55
    y = 2.35
    w = 2.35
    h = 2.75
    gap = 0.23

    col_centers = []

    for i, (short_title, full_title, node, fill, border) in enumerate(nodes):
        x = start_x + i * (w + gap)
        items = collect_child_names(node, max_items=4)

        if not items:
            items = ["No items generated"]

        add_map_column(
            slide,
            i + 1,
            short_title,
            items,
            x,
            y,
            w,
            h,
            fill,
            border,
        )

        col_centers.append((x + w / 2, y + h / 2))

    # 横向箭头，避免交叉
    for i in range(len(col_centers) - 1):
        x1 = start_x + i * (w + gap) + w
        x2 = start_x + (i + 1) * (w + gap)
        yy = y + h / 2

        add_connector(slide, x1, yy, x2, yy, COLOR_MUTED)

        arrow = slide.shapes.add_textbox(
            Inches(x1 + 0.03),
            Inches(yy - 0.17),
            Inches(0.22),
            Inches(0.25),
        )
        tf = arrow.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = "→"
        set_run_font(r, size=15, bold=True, color=COLOR_MUTED)

    # 底部一句总结，不再塞大 evidence 框
    add_card(
        slide,
        "How to read this map",
        "Start from prerequisite concepts, then understand the problem, method, experiments, and final takeaways.",
        1.1,
        5.55,
        11.1,
        0.62,
        fill=COLOR_WHITE,
        border=COLOR_GRAY,
        title_color=COLOR_DARK,
        body_size=11.5,
    )


def add_mindmap_evidence_footer(slide, evidence):
    """
    思维导图页只放很小的 evidence anchors，不放大段原文，避免压住图。
    """
    if not evidence:
        text = "Evidence anchors: generated from paragraph-level paper evidence."
    else:
        anchors = []
        for ev in evidence[:5]:
            pid = ev.get("paragraph_id", "?")
            page = ev.get("page_start", "?")
            section = ev.get("section_guess", "Other")
            anchors.append(f"P{pid}/Page {page}/{section}")
        text = "Evidence anchors: " + " · ".join(anchors)

    box = slide.shapes.add_textbox(
        Inches(0.75),
        Inches(6.52),
        Inches(11.8),
        Inches(0.28),
    )
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = truncate(text, 170)
    set_run_font(r, size=8.8, color=COLOR_MUTED)


# ---------- 特殊页面：方法流程图 ----------

def draw_method_flow(slide, method_flow):
    steps = method_flow[:5]
    if not steps:
        add_card(slide, "No method flow found", "The planner did not provide method steps.", 1, 2, 11, 2)
        return

    y = 2.15
    box_w = 2.15
    gap = 0.32
    start_x = 0.55

    for i, step in enumerate(steps):
        x = start_x + i * (box_w + gap)
        title = f"Step {step.get('step_no', i+1)}: {truncate(step.get('name', ''), 28)}"
        desc = truncate(step.get("description", ""), 95)

        add_card(
            slide,
            title,
            desc,
            x,
            y,
            box_w,
            1.45,
            fill=COLOR_WHITE,
            border=COLOR_BLUE,
            title_color=COLOR_BLUE,
            body_size=11,
        )

        if i < len(steps) - 1:
            add_connector(slide, x + box_w, y + 0.72, x + box_w + gap, y + 0.72, COLOR_BLUE)
            arrow = slide.shapes.add_textbox(Inches(x + box_w + 0.05), Inches(y + 0.48), Inches(0.25), Inches(0.3))
            tf = arrow.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = "→"
            set_run_font(r, size=18, bold=True, color=COLOR_BLUE)


def normalize_pid(x):
    """
    支持 19、"19"、"P19"、"PP19"，统一显示成 P19。
    """
    import re
    m = re.search(r"\d+", str(x))
    if not m:
        return str(x)
    return f"P{int(m.group(0))}"


def add_method_evidence(slide, method_flow):
    lines = []

    for step in method_flow[:5]:
        ids = step.get("evidence_paragraph_ids", [])
        if ids:
            normalized_ids = [normalize_pid(x) for x in ids]
            lines.append(f"{step.get('name', 'Step')}: {', '.join(normalized_ids)}")

    add_bullet_box(
        slide,
        "Method steps are grounded in these paragraphs",
        lines if lines else ["No paragraph ids provided."],
        1.0,
        4.35,
        11.3,
        1.55,
        max_item_len=120,
    )


# ---------- 页面生成 ----------

def create_title_slide(prs, grounded):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    paper_title = grounded.get("paper_title", "Unknown Paper")

    add_label(slide, "PaperBridge", 0.7, 0.55, 1.55, 0.36, COLOR_BLUE)

    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.35), Inches(11.6), Inches(1.25))
    tf = title_box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Beginner-Oriented Paper Understanding"
    set_run_font(r, size=32, bold=True, color=COLOR_DARK)

    sub_box = slide.shapes.add_textbox(Inches(1.25), Inches(2.65), Inches(10.8), Inches(1.0))
    tf = sub_box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = truncate(paper_title, 135)
    set_run_font(r, size=20, bold=True, color=COLOR_BLUE)

    add_card(
        slide,
        "System Output",
        "Learning path · Grounded slides · Mind map · Chinese narration · Lecture video",
        2.0,
        4.35,
        9.3,
        0.95,
        fill=COLOR_WHITE,
        border=COLOR_GRAY,
        title_color=COLOR_DARK,
        body_size=15,
    )

    add_card(
        slide,
        "Current Improvement",
        "Each slide is now grounded by paragraph-level evidence from the original paper.",
        2.0,
        5.55,
        9.3,
        0.75,
        fill=COLOR_LIGHT_BLUE,
        border=COLOR_BLUE,
        title_color=COLOR_BLUE,
        body_size=13,
    )

    return slide


def create_general_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    slide_no = slide_data.get("slide_no", "")
    title = slide_data.get("annotated_title", slide_data.get("title", f"Slide {slide_no}"))
    purpose = slide_data.get("purpose", "")

    add_title(slide, f"{slide_no}. {title}", purpose)

    main_points = slide_data.get(
        "annotated_main_points",
        slide_data.get("main_points", [])
    )
    evidence = slide_data.get("verified_evidence", [])
    visual_type = slide_data.get("visual_type", "")

    add_label(slide, f"Visual Type: {visual_type}", 0.65, 1.28, 2.1, 0.32, COLOR_GREEN)

    # 左侧：主要观点，拉高
    add_bullet_box(
        slide,
        "Main points",
        main_points[:5],
        0.65,
        1.75,
        6.05,
        4.95,
        max_item_len=115,
    )

    # 右侧：论文依据，拉高
    add_evidence_box(
        slide,
        evidence,
        7.05,
        1.75,
        5.6,
        4.95,
        max_items=4,
    )

    return slide


def create_mind_map_slide(prs, slide_data, mind_map):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_title(
        slide,
        "5. 论文理解学习地图",
        "A beginner-friendly learning map from prerequisites to takeaways."
    )

    draw_mind_map(slide, mind_map)

    evidence = slide_data.get("verified_evidence", [])
    add_mindmap_evidence_footer(slide, evidence)

    return slide


def create_method_flow_slide(prs, slide_data, method_flow):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_title(
        slide,
        "6. 核心方法流程图",
        "The method is represented as a grounded visual pipeline."
    )

    draw_method_flow(slide, method_flow)
    add_method_evidence(slide, method_flow)

    evidence = slide_data.get("verified_evidence", [])
    add_evidence_box(slide, evidence, 1.0, 6.15, 11.3, 0.85, max_items=3)

    return slide



def chunk_items(items, size):
    for i in range(0, len(items), size):
        yield items[i:i + size]


def create_glossary_slide(prs, glossary_items, page_index, total_pages):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_title(
        slide,
        f"{10 + page_index}. 专业名词注释表（{page_index}/{total_pages}）",
        "Terms marked with 〔1〕〔2〕... in previous slides are explained here."
    )

    if not glossary_items:
        add_card(
            slide,
            "No glossary generated",
            "Run python src/term_explainer.py before generating PPT.",
            1.0, 1.7, 11.2, 1.2,
            fill=COLOR_WHITE,
            border=COLOR_GRAY,
            title_color=COLOR_DARK,
            body_size=14,
        )
        return slide

    # 每页 10 个，左右各 5 个，卡片更高更清楚
    left_items = glossary_items[:5]
    right_items = glossary_items[5:10]

    def add_glossary_card(item, x, y):
        marker = item.get("marker", "")
        term = truncate(item.get("term", ""), 28)
        explanation = truncate(item.get("explanation", ""), 90)
        first_slide = item.get("first_slide_no", "?")

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(5.55),
            Inches(0.92),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_WHITE
        shape.line.color.rgb = RGBColor(203, 213, 225)
        shape.line.width = Pt(1)

        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_left = Inches(0.16)
        tf.margin_right = Inches(0.16)
        tf.margin_top = Inches(0.08)
        tf.margin_bottom = Inches(0.06)

        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{marker}  {term}"
        set_run_font(r, size=13.5, bold=True, color=COLOR_BLUE)

        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        r2 = p2.add_run()
        r2.text = explanation
        set_run_font(r2, size=9.6, color=COLOR_TEXT)

        p3 = tf.add_paragraph()
        p3.space_before = Pt(1)
        r3 = p3.add_run()
        r3.text = f"First appears: Slide {first_slide}"
        set_run_font(r3, size=8.4, color=COLOR_MUTED)

    y0 = 1.35
    gap = 1.02

    for i, item in enumerate(left_items):
        add_glossary_card(item, 0.75, y0 + i * gap)

    for i, item in enumerate(right_items):
        add_glossary_card(item, 6.95, y0 + i * gap)

    # 底部只放一行很轻的说明，不再用大卡片
    box = slide.shapes.add_textbox(
        Inches(0.9),
        Inches(6.72),
        Inches(11.5),
        Inches(0.25),
    )
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Use the numbered markers in previous slides to look up explanations here."
    set_run_font(r, size=8.5, color=COLOR_MUTED)

    return slide


def create_glossary_slides(prs, glossary):
    if not glossary:
        create_glossary_slide(prs, [], 1, 1)
        return

    chunks = list(chunk_items(glossary, 10))
    total = len(chunks)

    for i, items in enumerate(chunks, start=1):
        create_glossary_slide(prs, items, i, total)


def create_ppt(
    grounded_path: str = GROUNDED_PATH,
    output_path: str = OUTPUT_PPT,
):
    grounded = json.loads(Path(grounded_path).read_text(encoding="utf-8"))

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slides = grounded.get("slides", [])

    # Slide 1: Title
    create_title_slide(prs, grounded)

    # Slides 2-4
    for no in [2, 3, 4]:
        s = get_slide_by_no(slides, no)
        if s:
            create_general_slide(prs, s)

    # Slide 5: Mind map
    s5 = get_slide_by_no(slides, 5)
    create_mind_map_slide(prs, s5 or {}, grounded.get("mind_map", {}))

    # Slide 6: Method flow
    s6 = get_slide_by_no(slides, 6)
    create_method_flow_slide(prs, s6 or {}, grounded.get("method_flow", []))

    # Slides 7-10
    for no in [7, 8, 9, 10]:
        s = get_slide_by_no(slides, no)
        if s:
            create_general_slide(prs, s)

    # Glossary slides: explain all technical terms
    create_glossary_slides(prs, grounded.get("global_glossary", []))

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"Grounded PPT 已生成：{output_path}")


if __name__ == "__main__":
    create_ppt()
